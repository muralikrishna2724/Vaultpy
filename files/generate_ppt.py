"""
generate_ppt.py — Generates the VaultPy project presentation
Run with: python generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # completely blank
title_layout = prs.slide_layouts[0]   # title + subtitle

# ───────────────────────────────────────────────
# Helpers
# ───────────────────────────────────────────────

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=RGBColor(0x1A,0x1A,0x1A),
                align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txBox

def add_slide_header(slide, title_text, bg_color=RGBColor(0x1E,0x3A,0x5F)):
    """Add a colored header bar at the top of the slide."""
    # Background rectangle for header
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    # Title text on header
    add_textbox(slide, title_text,
                left=Inches(0.3), top=Inches(0.1),
                width=Inches(12.7), height=Inches(0.9),
                font_size=26, bold=True,
                color=RGBColor(0xFF,0xFF,0xFF),
                align=PP_ALIGN.LEFT)

def add_bullet_textbox(slide, items, left, top, width, height,
                        font_size=17, color=RGBColor(0x1A,0x1A,0x1A)):
    """Add a textbox with bullet-style items (• prefix)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color


# ───────────────────────────────────────────────
# Slide 1 — BLANK (as requested)
# ───────────────────────────────────────────────
slide1 = prs.slides.add_slide(blank_layout)


# ───────────────────────────────────────────────
# Slide 2 — Introduction
# ───────────────────────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)
add_slide_header(slide2, "Introduction")

intro_points = [
    "In today's digital world, people use dozens of apps and websites — each needing its own password.",
    "Reusing the same password everywhere is risky. One breach can expose all your accounts.",
    "VaultPy is a lightweight, offline password manager built for the desktop.",
    "It lets users securely store, retrieve, and manage all their passwords in one place.",
    "All data is encrypted locally — nothing ever goes to the cloud or any remote server.",
    "VaultPy is written in Python and packaged as a standalone desktop app using PyWebView + Flask.",
]
add_bullet_textbox(slide2, intro_points,
                   left=Inches(0.5), top=Inches(1.3),
                   width=Inches(12.3), height=Inches(5.8))


# ───────────────────────────────────────────────
# Slide 3 — Problem Statement
# ───────────────────────────────────────────────
slide3 = prs.slides.add_slide(blank_layout)
add_slide_header(slide3, "Problem Statement")

prob_points = [
    "People struggle to remember strong, unique passwords for every account they use.",
    "Most people fall back to weak or repeated passwords, making them easy targets for hackers.",
    "Cloud-based password managers store your data on their servers — a potential privacy risk.",
    "If the service shuts down or gets hacked, your passwords could be exposed or lost.",
    "There is no easy, open-source, offline tool that a regular user can set up and run locally.",
    "Goal: Build a simple, secure, fully offline password manager that runs on the user's own machine.",
]
add_bullet_textbox(slide3, prob_points,
                   left=Inches(0.5), top=Inches(1.3),
                   width=Inches(12.3), height=Inches(5.8))


# ───────────────────────────────────────────────
# Slide 4 — Existing Methodologies
# ───────────────────────────────────────────────
slide4 = prs.slides.add_slide(blank_layout)
add_slide_header(slide4, "Existing Methodologies")

existing_points = [
    "Cloud-Based Managers (e.g., LastPass, Dashlane): Store passwords on remote servers. Convenient but dependent on internet access and third-party security.",
    "Browser Built-in Managers (e.g., Chrome, Firefox): Easy to use but tied to a specific browser and synced through Google/Mozilla accounts.",
    "Hardware Tokens (e.g., YubiKey): Very secure but expensive and not practical for everyday users.",
    "Plain Text / Spreadsheet Storage: Still widely used. Extremely insecure — no encryption at all.",
    "Open-Source Tools (e.g., KeePass): Good security, but complex to set up and has an outdated interface.",
    "Limitation of all existing tools: Either too complex, too cloud-reliant, or not beginner-friendly.",
]
add_bullet_textbox(slide4, existing_points,
                   left=Inches(0.5), top=Inches(1.3),
                   width=Inches(12.3), height=Inches(5.8))


# ───────────────────────────────────────────────
# Slide 5 — Proposed Solution
# ───────────────────────────────────────────────
slide5 = prs.slides.add_slide(blank_layout)
add_slide_header(slide5, "Proposed Solution — VaultPy")

sol_points = [
    "VaultPy is a fully offline, encrypted password manager that runs as a native desktop window.",
    "A single master password is used to unlock the vault — the user only needs to remember one password.",
    "Passwords are stored in a local encrypted file using AES-256-GCM — a modern, military-grade cipher.",
    "The master password is never stored anywhere — it is used only to derive the encryption key using Argon2id.",
    "The app is built with Python (Flask) and opens in a native window using PyWebView — no browser needed.",
    "Simple, clean web interface makes it easy for any user to add, view, and delete passwords.",
    "Packaged as a .exe installer for Windows — no Python installation required for end users.",
]
add_bullet_textbox(slide5, sol_points,
                   left=Inches(0.5), top=Inches(1.3),
                   width=Inches(12.3), height=Inches(5.8))


# ───────────────────────────────────────────────
# Slide 6 — Software & Hardware Specifications
# ───────────────────────────────────────────────
slide6 = prs.slides.add_slide(blank_layout)
add_slide_header(slide6, "Software & Hardware Specifications")

# Software section label
add_textbox(slide6, "Software Requirements",
            left=Inches(0.5), top=Inches(1.3),
            width=Inches(5.8), height=Inches(0.5),
            font_size=18, bold=True, color=RGBColor(0x1E,0x3A,0x5F))

sw_items = [
    "Python 3.10+",
    "Flask 3.x — Web framework (backend)",
    "PyWebView 4.x — Native desktop window wrapper",
    "argon2-cffi — Argon2id key derivation",
    "cryptography — AES-256-GCM encryption",
    "PyInstaller — Packaging to .exe",
    "OS: Windows 10/11",
]
add_bullet_textbox(slide6, sw_items,
                   left=Inches(0.5), top=Inches(1.85),
                   width=Inches(5.8), height=Inches(4.5), font_size=16)

# Hardware section label
add_textbox(slide6, "Hardware Requirements",
            left=Inches(7.0), top=Inches(1.3),
            width=Inches(5.8), height=Inches(0.5),
            font_size=18, bold=True, color=RGBColor(0x1E,0x3A,0x5F))

hw_items = [
    "Processor: Any modern x86-64 CPU",
    "RAM: Minimum 512 MB (64 MB reserved for Argon2id key derivation)",
    "Storage: ~50 MB for the installed app + vault file",
    "Display: 800×600 or higher resolution",
    "No internet connection required",
]
add_bullet_textbox(slide6, hw_items,
                   left=Inches(7.0), top=Inches(1.85),
                   width=Inches(5.8), height=Inches(4.5), font_size=16)

# Divider line
line = slide6.shapes.add_shape(1, Inches(6.5), Inches(1.3), Inches(0.05), Inches(5.5))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
line.line.fill.background()


# ───────────────────────────────────────────────
# Slide 7 — Description of Modules
# ───────────────────────────────────────────────
slide7 = prs.slides.add_slide(blank_layout)
add_slide_header(slide7, "Description of Modules")

modules = [
    ("crypto.py", "Handles all cryptography — generates random salts, derives the encryption key from the master password using Argon2id, and encrypts/decrypts vault data with AES-256-GCM."),
    ("vault_manager.py", "Manages the vault file on disk. Handles creating, unlocking, locking, saving the vault and CRUD operations (add, update, delete, get) for password entries."),
    ("app.py", "The Flask web server, acting as the backend. Defines all URL routes — unlock, dashboard, add entry, delete entry, lock, and a password reveal API."),
    ("Templates (HTML)", "The frontend UI — unlock.html, dashboard.html, add_entry.html — rendered by Flask and displayed in the PyWebView native window."),
    ("build.bat / installer.iss", "Build scripts that package the entire app into a Windows .exe installer using PyInstaller and Inno Setup."),
]

y = Inches(1.35)
for mod_name, mod_desc in modules:
    add_textbox(slide7, mod_name,
                left=Inches(0.5), top=y,
                width=Inches(2.8), height=Inches(0.9),
                font_size=15, bold=True, color=RGBColor(0x1E,0x3A,0x5F))
    add_textbox(slide7, mod_desc,
                left=Inches(3.4), top=y,
                width=Inches(9.4), height=Inches(0.9),
                font_size=14, color=RGBColor(0x33,0x33,0x33))
    y += Inches(1.05)


# ───────────────────────────────────────────────
# Slide 8 — System Architecture (visual diagram)
# ───────────────────────────────────────────────
from pptx.enum.shapes import MSO_CONNECTOR_TYPE

slide8 = prs.slides.add_slide(blank_layout)
add_slide_header(slide8, "System Architecture")

# ── Color palette ──
C_USER    = RGBColor(0x27, 0xAE, 0x60)   # green  — User / UI
C_CRYPTO  = RGBColor(0x1E, 0x3A, 0x5F)   # navy   — crypto ops
C_KEY     = RGBColor(0xE6, 0x7E, 0x22)   # orange — key (sensitive)
C_STORE   = RGBColor(0x8E, 0x44, 0xAD)   # purple — storage
C_FLASK   = RGBColor(0x21, 0x8A, 0xC7)   # blue   — Flask / UI layer
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_ARROW   = RGBColor(0x55, 0x55, 0x55)

def draw_box(slide, label, sublabel, left, top, width, height, fill_color):
    """Draw a filled rounded rectangle with label text."""
    box = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE (MSO_SHAPE_TYPE 5)
        left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.width = Pt(0.5)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = C_WHITE

    if sublabel:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sublabel
        r2.font.size = Pt(9)
        r2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

def draw_arrow(slide, x1, y1, x2, y2):
    """Draw a straight connector arrow between two points."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2
    )
    conn.line.color.rgb = C_ARROW
    conn.line.width = Pt(1.5)

# ── Layout: two columns, left = write path, right = read path ──
# Shared top nodes
bw = Inches(2.4)
bh = Inches(0.58)

# Row y positions
r0 = Inches(1.25)
r1 = Inches(2.05)
r2 = Inches(2.85)
r3 = Inches(3.65)
r4 = Inches(4.45)
r5 = Inches(5.25)
r6 = Inches(6.05)

mid_x = Inches(5.47)   # centre x for shared boxes

# ── Shared nodes (centre column) ──
draw_box(slide8, "User", "Types master password",
         mid_x, r0, bw, bh, C_USER)

draw_box(slide8, "Argon2id Key Derivation", "salt + master password → 256-bit key",
         mid_x, r1, bw, bh, C_CRYPTO)

draw_box(slide8, "256-bit Encryption Key", "Lives in memory only — never saved",
         mid_x, r2, bw, bh, C_KEY)

# ── Left branch — Encrypt / Write ──
lx = Inches(1.5)
draw_box(slide8, "AES-256-GCM Encrypt", "Encrypts vault entries",
         lx, r3, bw, bh, C_CRYPTO)
draw_box(slide8, "vault.enc  (disk)", "Encrypted file on local disk",
         lx, r4, bw, bh, C_STORE)

# ── Right branch — Decrypt / Read ──
rx = Inches(9.43)
draw_box(slide8, "AES-256-GCM Decrypt", "Decrypts vault entries",
         rx, r3, bw, bh, C_CRYPTO)
draw_box(slide8, "In-Memory Entries", "Decrypted passwords in RAM",
         rx, r4, bw, bh, C_STORE)

# ── Bottom shared: Flask → PyWebView ──
draw_box(slide8, "Flask Web Server", "127.0.0.1:5000  (background thread)",
         mid_x, r5, bw, bh, C_FLASK)
draw_box(slide8, "PyWebView Desktop Window", "Native app window — no browser",
         mid_x, r6, bw, bh, C_USER)

# ─── Salt side-note box ───
draw_box(slide8, "Salt", "Stored in vault.enc",
         Inches(3.1), r1, Inches(1.6), bh, RGBColor(0x7F, 0x8C, 0x8D))

# ── Arrows ──
cx = mid_x + bw / 2          # centre-x of shared column boxes
mid_bottom = lambda row: row + bh   # bottom-centre y of a box

# User → Argon2id
draw_arrow(slide8, cx, mid_bottom(r0), cx, r1)
# Argon2id → Key
draw_arrow(slide8, cx, mid_bottom(r1), cx, r2)
# Salt → Argon2id
draw_arrow(slide8, Inches(3.1) + Inches(0.8), r1 + bh/2,
           mid_x, r1 + bh/2)
# Key → Encrypt (left branch)
draw_arrow(slide8, mid_x, r2 + bh/2, lx + bw, r2 + bh/2)
draw_arrow(slide8, lx + bw/2, r2 + bh/2, lx + bw/2, r3)
# Key → Decrypt (right branch)
draw_arrow(slide8, mid_x + bw, r2 + bh/2, rx, r2 + bh/2)
draw_arrow(slide8, rx + bw/2, r2 + bh/2, rx + bw/2, r3)
# Encrypt → vault.enc
draw_arrow(slide8, lx + bw/2, mid_bottom(r3), lx + bw/2, r4)
# Decrypt → In-memory
draw_arrow(slide8, rx + bw/2, mid_bottom(r3), rx + bw/2, r4)
# In-memory → Flask
draw_arrow(slide8, rx + bw/2, mid_bottom(r4), cx, r5 + bh/2)
draw_arrow(slide8, cx, r5 + bh/2, cx, r5)         # connect bend to Flask top
# Flask → PyWebView
draw_arrow(slide8, cx, mid_bottom(r5), cx, r6)


# ───────────────────────────────────────────────
# Slide 9 — Student Contributions
# ───────────────────────────────────────────────
slide9 = prs.slides.add_slide(blank_layout)
add_slide_header(slide9, "Contribution of the Students")

contrib_points = [
    "Designed and implemented the full cryptography layer — Argon2id key derivation + AES-256-GCM encryption (crypto.py).",
    "Built the VaultManager class for reading/writing the encrypted vault file including all CRUD operations (vault_manager.py).",
    "Developed the Flask backend with all routes — unlock, dashboard, add/delete entries, lock, and password reveal API (app.py).",
    "Designed the frontend HTML/CSS pages — unlock screen, password dashboard, and add-entry form.",
    "Set up the PyWebView integration to package the web UI as a native desktop window without a browser.",
    "Created the Windows build pipeline — PyInstaller spec, Inno Setup installer script, and build batch file.",
    "Conducted testing for encryption correctness, wrong-password handling, and session security.",
]
add_bullet_textbox(slide9, contrib_points,
                   left=Inches(0.5), top=Inches(1.3),
                   width=Inches(12.3), height=Inches(5.8))


# ───────────────────────────────────────────────
# Slide 10 — Results & Discussion
# ───────────────────────────────────────────────
slide10 = prs.slides.add_slide(blank_layout)
add_slide_header(slide10, "Results & Discussion")

results_points = [
    "The vault file is fully encrypted — opening vault.enc in a text editor shows only unreadable ciphertext.",
    "Entering a wrong master password returns a clean error message and never decrypts anything.",
    "Argon2id key derivation takes ~0.5–1 second intentionally — this slows down brute-force attacks significantly.",
    "AES-256-GCM authentication tag ensures any file tampering is detected immediately on next unlock.",
    "The app launches in a native window — no browser address bar, no URL visible to the user.",
    "Passwords can be added, viewed (with reveal button), and deleted — all without leaving the local machine.",
    "The .exe installer works on Windows 10/11 without requiring Python to be pre-installed.",
    "Minor limitation: Currently no password generator or import/export feature — planned for future versions.",
]
add_bullet_textbox(slide10, results_points,
                   left=Inches(0.5), top=Inches(1.3),
                   width=Inches(12.3), height=Inches(5.8))


# ───────────────────────────────────────────────
# Slide 11 — Conclusion
# ───────────────────────────────────────────────
slide11 = prs.slides.add_slide(blank_layout)
add_slide_header(slide11, "Conclusion")

conc_points = [
    "VaultPy successfully achieves its goal — a simple, secure, and fully offline password manager.",
    "Users only need to remember one master password; everything else is safely encrypted on their own machine.",
    "The use of Argon2id + AES-256-GCM ensures the vault is protected by modern, industry-standard cryptography.",
    "By avoiding cloud storage entirely, VaultPy removes a critical single point of failure found in most popular tools.",
    "The PyWebView-based packaging gives users a smooth desktop experience without needing any technical setup.",
    "Future scope includes: built-in password generator, CSV import/export, password strength checker, and auto-lock on idle.",
    "This project demonstrates how strong cryptography can be made practical and accessible for everyday users.",
]
add_bullet_textbox(slide11, conc_points,
                   left=Inches(0.5), top=Inches(1.3),
                   width=Inches(12.3), height=Inches(5.8))


# ───────────────────────────────────────────────
# Slide 12 — References
# ───────────────────────────────────────────────
slide12 = prs.slides.add_slide(blank_layout)
add_slide_header(slide12, "References")

refs = [
    "[1] Biryukov, A., Dinu, D., & Khovratovich, D. (2016). Argon2: The Memory-Hard Function for Password Hashing. IETF RFC 9106.",
    "[2] NIST SP 800-38D — Recommendation for Block Cipher Modes of Operation: Galois/Counter Mode (GCM). National Institute of Standards and Technology.",
    "[3] Flask Documentation. Pallets Projects. https://flask.palletsprojects.com/",
    "[4] pywebview Documentation. Roman Sirokov. https://pywebview.flowrl.com/",
    "[5] Python cryptography library — cryptography.io. https://cryptography.io/en/latest/",
    "[6] argon2-cffi Documentation. Hynek Schlawack. https://argon2-cffi.readthedocs.io/",
    "[7] PyInstaller Documentation. https://pyinstaller.org/en/stable/",
    "[8] Inno Setup Documentation. Jordan Russell. https://jrsoftware.org/ishelp/",
]
add_bullet_textbox(slide12, refs,
                   left=Inches(0.4), top=Inches(1.3),
                   width=Inches(12.5), height=Inches(5.8),
                   font_size=14, color=RGBColor(0x22,0x22,0x22))


# ───────────────────────────────────────────────
# Save
# ───────────────────────────────────────────────
out_path = r"c:\Users\surya\OneDrive\Desktop\app\files\VaultPy_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
